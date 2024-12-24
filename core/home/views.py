import os
from django.shortcuts import render
from django.http import HttpResponse
from pdf2docx import Converter
from pptx import Presentation
import fitz  
from PIL import Image
import zipfile
from django.shortcuts import render
from django.http import HttpResponse
from django.core.files.storage import default_storage
from django.contrib.auth.decorators import login_required
import zipfile
from datetime import datetime, timedelta
from django.shortcuts import render
from django.core.files.storage import default_storage
import pandas as pd
from pdfplumber import PDF
from django.http import HttpResponse
import pythoncom
import win32com.client
from django.http import HttpResponse
from django.core.files.storage import default_storage
from django.shortcuts import render
from fpdf import FPDF
from PyPDF2 import PdfReader
from django.core.files.storage import default_storage
import pdfplumber


def check_usage_limit(request, feature_name):
    if request.user.is_authenticated:
        return True  

    usage_count = request.session.get('tool_usage', {})
    last_reset = request.session.get('tool_usage_timestamp')

    if last_reset:
        last_reset_time = datetime.strptime(last_reset, "%Y-%m-%d %H:%M:%S")
        if datetime.now() - last_reset_time > timedelta(hours=24):
            usage_count = {}
            request.session['tool_usage'] = usage_count
            request.session['tool_usage_timestamp'] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    else:
        request.session['tool_usage_timestamp'] = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    count = usage_count.get(feature_name, 0)

    if count >= 2:
        return False 
    
    usage_count[feature_name] = count + 1
    request.session['tool_usage'] = usage_count
    return True

def index(request):
    return render(request, 'index.html')

def handle_upload(file):
    file_path = default_storage.save(file.name, file)
    return default_storage.path(file_path)

def pdf_to_word(request):
    feature_name = "PDF to Word"
    if request.method == "POST":
        if not check_usage_limit(request, feature_name):
            return render(request, 'popup.html', {'message': 'Daily Usage limit reached. Please log in to continue .'})

        file = request.FILES['file']
        pdf_path = handle_upload(file)
        docx_path = pdf_path.replace('.pdf', '.docx')
        converter = Converter(pdf_path)
        converter.convert(docx_path)
        converter.close()
        with open(docx_path, 'rb') as docx_file:
            response = HttpResponse(docx_file.read(), content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
            response['Content-Disposition'] = f'attachment; filename={os.path.basename(docx_path)}'
        os.remove(pdf_path)
        os.remove(docx_path)
        return response
    return render(request, 'pdf_to_word.html', {'feature': feature_name})


def pdf_to_ppt(request):
    feature_name = "PDF to PPT"
    if request.method == "POST":
        if not check_usage_limit(request, feature_name):
            return render(request, 'popup.html', {'message': 'Daily Usage limit reached. Please log in to continue.'})
        
        file = request.FILES['file']
        pdf_path = handle_upload(file)
        ppt_path = pdf_path.replace('.pdf', '.pptx')
        pdf_document = fitz.open(pdf_path)
        presentation = Presentation()
        
        for page_number in range(len(pdf_document)):
            slide = presentation.slides.add_slide(presentation.slide_layouts[6])
            page = pdf_document[page_number]
            pix = page.get_pixmap()
            image_path = f"{pdf_path}_{page_number}.png"
            pix.save(image_path)
            
            slide.shapes.add_picture(image_path, 0, 0, width=presentation.slide_width)
            os.remove(image_path)
        pdf_document.close()
        
        presentation.save(ppt_path)
        
        with open(ppt_path, 'rb') as ppt_file:
            response = HttpResponse(ppt_file.read(), content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
            response['Content-Disposition'] = f'attachment; filename={os.path.basename(ppt_path)}'
        os.remove(pdf_path)
        os.remove(ppt_path)
        return response
    
    return render(request, 'pdf_to_ppt.html', {'feature': 'PDF to PPT'})

def pdf_to_jpg(request):
    
    feature_name = "PDF to JPG"
    if request.method == "POST":
        
        if not check_usage_limit(request, feature_name):
            return render(request, 'popup.html', {'message': 'Usage limit reached. Please log in to continue.'})

        file = request.FILES['file']
        pdf_path = handle_upload(file)
        pdf_document = fitz.open(pdf_path)
        output_paths = []
        
        for page_number in range(len(pdf_document)):
            page = pdf_document[page_number]
            pix = page.get_pixmap()
            image_path = f"{pdf_path}_{page_number}.jpg"
            pix.save(image_path)
            output_paths.append(image_path)
            
        pdf_document.close()
        
        response = HttpResponse(content_type='application/zip')
        response['Content-Disposition'] = 'attachment; filename="converted_images.zip"'
        
        with zipfile.ZipFile(response, 'w') as zipf:
            
            for image_path in output_paths:
                zipf.write(image_path, os.path.basename(image_path))
                os.remove(image_path)
        os.remove(pdf_path)
        return response
    
    return render(request, 'pdf_to_jpg.html', {'feature': 'PDF to JPG'})

from PyPDF2 import PdfReader, PdfWriter

def lock_pdf(request):
    feature_name = "Lock PDF"
    if request.method == "POST":
        
        if not check_usage_limit(request, feature_name):
            return render(request, 'popup.html', {'message': 'Usage limit reached. Please log in to continue.'})
        
    if request.method == "POST":
        file = request.FILES['file']
        password = request.POST['password']
        pdf_path = handle_upload(file)
        locked_pdf_path = pdf_path.replace('.pdf', '_locked.pdf')
        
        reader = PdfReader(pdf_path)
        writer = PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
        writer.encrypt(password)
        
        with open(locked_pdf_path, 'wb') as locked_pdf:
            writer.write(locked_pdf)
        
        with open(locked_pdf_path, 'rb') as pdf_file:
            response = HttpResponse(pdf_file.read(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename={os.path.basename(locked_pdf_path)}'
        
        os.remove(pdf_path)
        os.remove(locked_pdf_path)
        return response
    return render(request, 'lock_pdf.html', {'feature': 'Lock PDF'})


def unlock_pdf(request):
    feature_name = "Unlock PDF"
    if request.method == "POST":
        
        if not check_usage_limit(request, feature_name):
            return render(request, 'popup.html', {'message': 'Usage limit reached. Please log in to continue.'})
        
        file = request.FILES['file']
        password = request.POST['password']
        pdf_path = handle_upload(file)
        unlocked_pdf_path = pdf_path.replace('.pdf', '_unlocked.pdf')
        
        reader = PdfReader(pdf_path)
        writer = PdfWriter()
        if reader.is_encrypted:
            reader.decrypt(password)
        
        for page in reader.pages:
            writer.add_page(page)
        
        with open(unlocked_pdf_path, 'wb') as unlocked_pdf:
            writer.write(unlocked_pdf)
        
        with open(unlocked_pdf_path, 'rb') as pdf_file:
            response = HttpResponse(pdf_file.read(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename={os.path.basename(unlocked_pdf_path)}'
        
        os.remove(pdf_path)
        os.remove(unlocked_pdf_path)
        return response
    return render(request, 'unlock_pdf.html', {'feature': 'Unlock PDF'})


from PIL import Image

def jpg_to_pdf(request):
    feature_name = "JPG to PDF"
    if request.method == "POST":
        
        if not check_usage_limit(request, feature_name):
            return render(request, 'popup.html', {'message': 'Usage limit reached. Please log in to continue.'})
        files = request.FILES.getlist('file')
        image_paths = [handle_upload(f) for f in files]
        pdf_path = "output.pdf"
        
        image_list = []
        for img_path in image_paths:
            img = Image.open(img_path)
            if img.mode != 'RGB':
                img = img.convert('RGB')
            image_list.append(img)
        
        image_list[0].save(pdf_path, save_all=True, append_images=image_list[1:])
        
        with open(pdf_path, 'rb') as pdf_file:
            response = HttpResponse(pdf_file.read(), content_type='application/pdf')
            response['Content-Disposition'] = 'attachment; filename="converted.pdf"'
        
        for img_path in image_paths:
            os.remove(img_path)
        os.remove(pdf_path)
        return response
    return render(request, 'jpg_to_pdf.html', {'feature': 'JPG to PDF'})


import comtypes.client
import comtypes.client
import pythoncom  

def ppt_to_pdf(request):
    feature_name = "PPT to PDF"
    if request.method == "POST":
        
        if not check_usage_limit(request, feature_name):
            return render(request, 'popup.html', {'message': 'Usage limit reached. Please log in to continue.'})
        
        file = request.FILES['file']
        ppt_path = handle_upload(file)
        pdf_path = ppt_path.replace('.pptx', '.pdf')

        try:
            pythoncom.CoInitialize()  
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            presentation = powerpoint.Presentations.Open(ppt_path, WithWindow=False)
            presentation.SaveAs(pdf_path, 32)  
            presentation.Close()
            powerpoint.Quit()
        finally:
            pythoncom.CoUninitialize()  

        with open(pdf_path, 'rb') as pdf_file:
            response = HttpResponse(pdf_file.read(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename={os.path.basename(pdf_path)}'

        os.remove(ppt_path)
        os.remove(pdf_path)
        return response

    return render(request, 'ppt_to_pdf.html', {'feature': 'PPT to PDF'})



from django.core.files.storage import default_storage
from django.http import HttpResponse
import comtypes.client
import pythoncom  

def word_to_pdf(request):
    feature_name = "Word to PDF"
    if request.method == "POST":
        
        if not check_usage_limit(request, feature_name):
            return render(request, 'popup.html', {'message': 'Usage limit reached. Please log in to continue.'})
        
        pythoncom.CoInitialize()
        try:
            file = request.FILES['file']
            word_path = handle_upload(file)
            pdf_path = word_path.replace('.docx', '.pdf')
            
            word = comtypes.client.CreateObject('Word.Application')
            doc = word.Documents.Open(word_path)
            doc.SaveAs(pdf_path, FileFormat=17) 
            doc.Close()
            word.Quit()
            
            with open(pdf_path, 'rb') as pdf_file:
                response = HttpResponse(pdf_file.read(), content_type='application/pdf')
                response['Content-Disposition'] = f'attachment; filename={os.path.basename(pdf_path)}'
            
            os.remove(word_path)
            os.remove(pdf_path)
            return response
        finally:
            
            pythoncom.CoUninitialize()

    return render(request, 'word_to_pdf.html', {'feature': 'Word to PDF'})




from PyPDF2 import PdfMerger

def merge_pdf(request):
    feature_name = "Merge PDF"
    if request.method == "POST":
        
        if not check_usage_limit(request, feature_name):
            return render(request, 'popup.html', {'message': 'Usage limit reached. Please log in to continue.'})
        
        files = request.FILES.getlist('file')
        pdf_paths = [handle_upload(file) for file in files]
        
        merger = PdfMerger()
        for pdf in pdf_paths:
            merger.append(pdf)
        
        merged_pdf_path = 'merged_output.pdf'
        merger.write(merged_pdf_path)
        merger.close()
        
        with open(merged_pdf_path, 'rb') as pdf_file:
            response = HttpResponse(pdf_file.read(), content_type='application/pdf')
            response['Content-Disposition'] = 'attachment; filename="merged_output.pdf"'
        
        for pdf_path in pdf_paths:
            os.remove(pdf_path)
        os.remove(merged_pdf_path)
        return response
    return render(request, 'merge_pdf.html', {'feature': 'Merge PDF'})



def split_pdf(request):
    feature_name = "Split PDF"
    if request.method == "POST":
        if not check_usage_limit(request, feature_name):
            return render(request, 'popup.html', {'message': 'Daily Usage limit reached. Please log in to continue.'})

        file = request.FILES['file']
        pdf_path = handle_upload(file)
        pdf_document = fitz.open(pdf_path)

        output_files = []
        for page_num in range(len(pdf_document)):
            output_path = pdf_path.replace('.pdf', f'_page_{page_num + 1}.pdf')
            pdf_writer = fitz.open()
            pdf_writer.insert_pdf(pdf_document, from_page=page_num, to_page=page_num)
            pdf_writer.save(output_path)
            pdf_writer.close()
            output_files.append(output_path)

        pdf_document.close()

        zip_path = pdf_path.replace('.pdf', '_split_pages.zip')
        with zipfile.ZipFile(zip_path, 'w') as zipf:
            for file_path in output_files:
                zipf.write(file_path, os.path.basename(file_path))
                os.remove(file_path)

        with open(zip_path, 'rb') as zip_file:
            response = HttpResponse(zip_file.read(), content_type='application/zip')
            response['Content-Disposition'] = f'attachment; filename={os.path.basename(zip_path)}'
        
        os.remove(pdf_path)
        os.remove(zip_path)
        return response
    return render(request, 'pdf_split.html', {'feature': feature_name})


def pdf_to_excel(request):
    feature_name = "PDF to Excel"
    if request.method == "POST":
        if not check_usage_limit(request, feature_name):
            return render(request, 'popup.html', {'message': 'Daily Usage limit reached. Please log in to continue.'})
        
        file = request.FILES['file']
        pdf_path = handle_upload(file)
        excel_path = pdf_path.replace('.pdf', '.xlsx')

        try:
            with PDF(open(pdf_path, 'rb')) as pdf:
                all_data = []
                for page in pdf.pages:
                    tables = page.extract_tables()
                    for table in tables:
                        all_data.extend(table)
            
            df = pd.DataFrame(all_data)
            df.to_excel(excel_path, index=False, header=False)

            with open(excel_path, 'rb') as excel_file:
                response = HttpResponse(excel_file.read(), content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
                response['Content-Disposition'] = f'attachment; filename={os.path.basename(excel_path)}'

            os.remove(pdf_path)
            os.remove(excel_path)
            return response
        except Exception as e:
            os.remove(pdf_path)
            return render(request, 'error.html', {'message': f"An error occurred: {e}"})
    return render(request, 'pdf_to_excel.html', {'feature': feature_name})


def excel_to_pdf(request):
    feature_name = "Excel to PDF"
    if request.method == "POST":

        if not check_usage_limit(request, feature_name):
            return render(request, 'popup.html', {'message': 'Usage limit reached. Please log in to continue.'})

        file = request.FILES['file']
        excel_path = handle_upload(file)
        pdf_path = excel_path.replace('.xlsx', '.pdf').replace('.xls', '.pdf')

        try:
            pythoncom.CoInitialize()  
            excel = win32com.client.Dispatch("Excel.Application")
            workbook = excel.Workbooks.Open(excel_path)
            workbook.SaveAs(pdf_path, FileFormat=57)  # 57 is pfd formta
            workbook.Close()
            excel.Quit()
        except Exception as e:
            return render(request, 'popup.html', {'message': f'Error: {str(e)}'})
        finally:
            pythoncom.CoUninitialize()  

        with open(pdf_path, 'rb') as pdf_file:
            response = HttpResponse(pdf_file.read(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename={os.path.basename(pdf_path)}'

        
        os.remove(excel_path)
        os.remove(pdf_path)

        return response

    return render(request, 'excel_to_pdf.html', {'feature': feature_name})


def html_to_pdf(request):
    feature_name = "HTML to PDF"
    if request.method == "POST":
        if not check_usage_limit(request, feature_name):
            return render(request, 'popup.html', {'message': 'Usage limit reached. Please log in to continue.'})
        
        content = request.POST.get('html_content', '')
        if not content.strip():
            return render(request, 'html_to_pdf.html', {'feature': feature_name, 'error': 'No content provided.'})
        
        try:
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)

            for line in content.splitlines():
                pdf.multi_cell(0, 10, line)

            pdf_path = "output.pdf"
            pdf.output(pdf_path)

            with open(pdf_path, 'rb') as pdf_file:
                response = HttpResponse(pdf_file.read(), content_type='application/pdf')
                response['Content-Disposition'] = 'attachment; filename="converted.pdf"'
            
            return response
        
        finally:
            if os.path.exists(pdf_path):
                os.remove(pdf_path)
    
    return render(request, 'html_to_pdf.html', {'feature': feature_name})


def pdf_to_html(request):
    feature_name = "PDF to HTML"
    if request.method == "POST":
        if not check_usage_limit(request, feature_name):
            return render(request, 'popup.html', {'message': 'Daily Usage limit reached. Please log in to continue.'})

        file = request.FILES.get('file')
        if not file:
            return render(request, 'pdf_to_html.html', {'feature': feature_name, 'error': 'Please upload a PDF file.'})
        
        pdf_path = handle_upload(file)
        html_path = pdf_path.replace('.pdf', '.html')

        try:
            html_content = "<html><body>"
            with pdfplumber.open(pdf_path) as pdf:
                for page in pdf.pages:
                    text = page.extract_text()
                    if text:
                        html_content += f"<h3>Page {page.page_number}</h3>"
                        html_content += f"<p>{text.replace('\n', '<br>')}</p><hr>"
                    else:
                        html_content += f"<h3>Page {page.page_number}</h3><p>No text content available</p><hr>"
            html_content += "</body></html>"
            
            with open(html_path, 'w', encoding='utf-8') as html_file:
                html_file.write(html_content)

            with open(html_path, 'rb') as html_file:
                response = HttpResponse(html_file.read(), content_type='text/html')
                response['Content-Disposition'] = f'attachment; filename={os.path.basename(html_path)}'
            
            return response
        finally:
            if os.path.exists(pdf_path):
                os.remove(pdf_path)
            if os.path.exists(html_path):
                os.remove(html_path)
    
    return render(request, 'pdf_to_html.html', {'feature': feature_name})


import pikepdf
from PyPDF2 import PdfReader, PdfWriter

def compress_pdf(request):
    feature_name = "Compress PDF"
    if request.method == "POST":
        if not check_usage_limit(request, feature_name):
            return render(request, 'popup.html', {'message': 'Daily Usage limit reached. Please log in to continue.'})

        file = request.FILES['file']
        pdf_path = handle_upload(file)
        compressed_pdf_path = pdf_path.replace('.pdf', '_compressed.pdf')
        
        with pikepdf.open(pdf_path) as pdf:
            pdf.save(compressed_pdf_path, compress_streams=True)

        with open(compressed_pdf_path, 'rb') as compressed_file:
            response = HttpResponse(compressed_file.read(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename={os.path.basename(compressed_pdf_path)}'

        os.remove(pdf_path)
        os.remove(compressed_pdf_path)
        return response

    return render(request, 'compress_pdf.html', {'feature': feature_name})


from PyPDF2 import PdfReader, PdfWriter
from io import BytesIO
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import math

def add_watermark(request):
    feature_name = "Add Watermark"
    if request.method == "POST":
        if not check_usage_limit(request, feature_name):
            return render(request, 'popup.html', {'message': 'Daily Usage limit reached. Please log in to continue.'})

        file = request.FILES['file']
        watermark_text = request.POST.get('watermark_text', '')  
        if not watermark_text:
            return render(request, 'popup.html', {'message': 'Please enter a watermark text.'})

        pdf_path = handle_upload(file)
        watermarked_pdf_path = pdf_path.replace('.pdf', '_watermarked.pdf')

        packet = BytesIO()
        can = canvas.Canvas(packet, pagesize=letter)
        can.setFont("Helvetica", 36)
        can.setFillColorRGB(0.5, 0.5, 0.5, 0.5)  

        width, height = letter
        x_center = width / 2
        y_center = height / 2

       
        can.translate(x_center, y_center)  
        can.rotate(45)  
        can.drawString(-100, 0, watermark_text)  

        can.save()

        packet.seek(0)
        watermark_pdf = PdfReader(packet)
        watermark_page = watermark_pdf.pages[0]

        pdf_reader = PdfReader(pdf_path)
        pdf_writer = PdfWriter()

        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            page.merge_page(watermark_page)
            pdf_writer.add_page(page)

        with open(watermarked_pdf_path, 'wb') as output_pdf:
            pdf_writer.write(output_pdf)

        with open(watermarked_pdf_path, 'rb') as watermarked_file:
            response = HttpResponse(watermarked_file.read(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename={os.path.basename(watermarked_pdf_path)}'

        os.remove(pdf_path)
        os.remove(watermarked_pdf_path)
        return response

    return render(request, 'add_watermark.html', {'feature': feature_name})


from PyPDF2 import PdfReader, PdfWriter

def rotate_pdf(request):
    feature_name = "Rotate PDF"
    if request.method == "POST":
        if not check_usage_limit(request, feature_name):
            return render(request, 'popup.html', {'message': 'Daily Usage limit reached. Please log in to continue.'})

        file = request.FILES['file']
        rotation_angle = int(request.POST.get('angle', 90))  
        pdf_path = handle_upload(file)
        rotated_pdf_path = pdf_path.replace('.pdf', '_rotated.pdf')

        pdf_reader = PdfReader(pdf_path)
        pdf_writer = PdfWriter()

        for page_num in range(len(pdf_reader.pages)):
            page = pdf_reader.pages[page_num]
            page.rotate(rotation_angle)  
            pdf_writer.add_page(page)

        with open(rotated_pdf_path, 'wb') as output_pdf:
            pdf_writer.write(output_pdf)

        with open(rotated_pdf_path, 'rb') as rotated_file:
            response = HttpResponse(rotated_file.read(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename={os.path.basename(rotated_pdf_path)}'

        os.remove(pdf_path)
        os.remove(rotated_pdf_path)
        return response

    return render(request, 'rotate_pdf.html', {'feature': feature_name})

from PyPDF2 import PdfReader, PdfWriter
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import io

def add_numbering_pdf(request):
    if request.method == "POST":
        file = request.FILES['file']
        pdf_path = handle_upload(file)
        output_pdf_path = pdf_path.replace('.pdf', '_numbered.pdf')
    
        reader = PdfReader(pdf_path)
        writer = PdfWriter()
        packet = io.BytesIO()
        c = canvas.Canvas(packet, pagesize=letter)
        for page_num in range(len(reader.pages)):
            c.setFont("Helvetica", 12)
            c.drawString(10, 780, str(page_num + 1)) 
            c.showPage()
        c.save()

        packet.seek(0)
        number_pdf = PdfReader(packet)
        for i in range(len(reader.pages)):
            page = reader.pages[i]
            page.merge_page(number_pdf.pages[i])
            writer.add_page(page)
        
        with open(output_pdf_path, "wb") as output_pdf:
            writer.write(output_pdf)


        with open(output_pdf_path, 'rb') as numbered_pdf:
            response = HttpResponse(numbered_pdf.read(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename={os.path.basename(output_pdf_path)}'
        os.remove(pdf_path)
        os.remove(output_pdf_path)
        return response
    return render(request, 'add_numbering_pdf.html')

from PyPDF2.generic import NameObject, IndirectObject
from PyPDF2 import PdfReader, PdfWriter
from reportlab.pdfgen import canvas
import io

def sign_pdf(request):
    if request.method == "POST":
        file = request.FILES['file']
        pdf_path = handle_upload(file)
        output_pdf_path = pdf_path.replace('.pdf', '_signed.pdf')
        
        packet = io.BytesIO()
        c = canvas.Canvas(packet)
        c.setFont("Helvetica", 12)
        c.drawString(500, 30, "Your_Name")
        c.showPage()
        c.save()

        
        packet.seek(0)
        signature_pdf = PdfReader(packet)
        reader = PdfReader(pdf_path)
        writer = PdfWriter()
        
        for i in range(len(reader.pages)):
            page = reader.pages[i]
            page.merge_page(signature_pdf.pages[0])
            writer.add_page(page)

        with open(output_pdf_path, "wb") as output_pdf:
            writer.write(output_pdf)

        with open(output_pdf_path, 'rb') as signed_pdf:
            response = HttpResponse(signed_pdf.read(), content_type='application/pdf')
            response['Content-Disposition'] = f'attachment; filename={os.path.basename(output_pdf_path)}'
        os.remove(pdf_path)
        os.remove(output_pdf_path)
        return response
    return render(request, 'sign_pdf.html')

import pytesseract
from pdf2image import convert_from_path
import os

def ocr_pdf(request):
    if request.method == "POST":
        file = request.FILES['file']
        pdf_path = handle_upload(file)
        output_txt_path = pdf_path.replace('.pdf', '_ocr.txt')
        
        images = convert_from_path(pdf_path)
    
        with open(output_txt_path, "w") as output_txt:
            for page_num, image in enumerate(images):
                text = pytesseract.image_to_string(image)
                output_txt.write(f"Page {page_num + 1}:\n{text}\n\n")

        with open(output_txt_path, 'r') as ocr_file:
            response = HttpResponse(ocr_file.read(), content_type='text/plain')
            response['Content-Disposition'] = f'attachment; filename={os.path.basename(output_txt_path)}'
        os.remove(pdf_path)
        os.remove(output_txt_path)
        return response
    return render(request, 'ocr_pdf.html')


from django.shortcuts import render, redirect
from django.contrib.auth.models import User
from django.contrib.auth import authenticate, login, logout
from django.contrib import messages
from django.contrib.auth.decorators import login_required
import re  

def register_view(request):
    if request.method == 'POST':
        username = request.POST['username'].strip()
        email = request.POST['email'].strip()
        password = request.POST['password']
        confirm_password = request.POST['confirm_password']

        if not username or not email or not password or not confirm_password:
            messages.error(request, "All fields are required.")
            return redirect('register')
        
        if len(username) < 3:
            messages.error(request, "Username must be at least 3 characters long.")
            return redirect('register')

        email_regex = r'^[a-zA-Z0-9_.+-]+@[a-zA-Z0-9-]+\.[a-zA-Z0-9-.]+$'
        if not re.match(email_regex, email):
            messages.error(request, "Invalid email format.")
            return redirect('register')

        
        if password != confirm_password:
            messages.error(request, "Passwords do not match.")
            return redirect('register')

        if len(password) < 8:
            messages.error(request, "Password must be at least 8 characters long.")
            return redirect('register')

        if User.objects.filter(username=username).exists():
            messages.error(request, "Username already exists.")
            return redirect('register')

        if User.objects.filter(email=email).exists():
            messages.error(request, "Email already exists.")
            return redirect('register')

        user = User.objects.create_user(username=username, email=email, password=password)
        user.save()
        messages.success(request, "Registration successful! Please login.")
        return redirect('login')

    return render(request, 'register.html')


def login_view(request):
    if request.method == 'POST':
        username = request.POST['username'].strip()
        password = request.POST['password']

        if not username or not password:
            messages.error(request, "Both username and password are required.")
            return redirect('login')

        user = authenticate(request, username=username, password=password)
        if user is not None:
            login(request, user)
            return redirect('index')
        else:
            messages.error(request, "Invalid username or password.")
            return redirect('login')

    return render(request, 'login.html')


@login_required
def profile_view(request):
    return render(request, 'profile.html')


@login_required
def update_view(request):
    if request.method == 'POST':
        username = request.POST['username'].strip()
        email = request.POST['email'].strip()

        if not username or not email:
            messages.error(request, "All fields are required.")
            return redirect('update')

        user = request.user


        if user.username != username and User.objects.filter(username=username).exists():
            messages.error(request, "Username already exists.")
            return redirect('update')

        if user.email != email and User.objects.filter(email=email).exists():
            messages.error(request, "Email already exists.")
            return redirect('update')

        user.username = username
        user.email = email
        user.save()

        messages.success(request, "Account updated successfully.")
        return redirect('profile')

    return render(request, 'update.html')

@login_required
def delete_view(request):
    if request.method == 'POST':
        request.user.delete()
        messages.success(request, "Account deleted successfully.")
        return redirect('index')
    return render(request, 'delete.html')


@login_required
def logout_view(request):
    logout(request)
    return redirect('index')


